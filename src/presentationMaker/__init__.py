from pptx import *
from pptx.dml.color import RGBColor
from pptx.util import Pt
from overload import overload
import PIL
import os

typeToLayout = {
    "title": 0,
    "bullets": 1,
    "content": 2
}


def hex_to_rgb(value):
    value = value.lstrip('#')
    lv = len(value)
    r, g, b = tuple(int(value[i:i + lv // 3], 16) for i in range(0, lv, lv // 3))
    return RGBColor(r, g, b)


class TextData:
    @overload
    def __init__(self, data):
        if isinstance(data, str):
            self.text = data
            self.color = None
            self.font = None
            self.fontSize = None
        else:
            self.text = data["text"]
            self.color = data.get("color", None)
            self.font = data.get("font", None)
            self.fontSize = data.get("fontSize", None)

    @__init__.add
    def __init__(self, data, color, font, fontSize):
        self.__init__(data)
        if not self.color: self.color = color
        if not self.font: self.font = font
        if not self.fontSize: self.fontSize = fontSize

    def applyTo(self, font):
        if self.color:
            font.color.rgb = hex_to_rgb(self.color)
        if self.font:
            font.name = self.font
        if self.fontSize:
            font.size = Pt(int(self.fontSize))


class BulletData:
    def __init__(self, data):
        self.bullets = []

        textInfo = data["text"]
        color = None
        font = None
        fontSize = None

        if isinstance(textInfo, str):
            file = open(textInfo, "r")
            self.fromFile(file, color, font, fontSize)
        else:
            color = textInfo.get("color", None)
            font = textInfo.get("font", None)
            fontSize = textInfo.get("fontSize", None)
            if "paragraphs" in textInfo:
                for paragraph in textInfo["paragraphs"]:
                    self.bullets.append(TextData(paragraph, color, font, fontSize))
            else:
                file = open(textInfo["file"], "r")
                self.fromFile(file, color, font, fontSize)

    def fromFile(self, file, color, font, fontSize):
        text = ""
        lines = file.readlines()
        for i in range(0, len(lines)):
            line = lines[i]
            text += line.rstrip()
            text += " "
            if line.isspace() or i == len(lines) - 1:
                self.bullets.append(TextData(text, color, font, fontSize))
                text = ""


def makePresentation(description):
    slides = description["slides"]
    runPath = os.path.dirname(os.path.realpath(__file__))
    prs = Presentation(os.path.join(runPath, "template.pptx"))
    for slideInfo in slides:
        slideType = slideInfo["type"]
        layout = prs.slide_layouts[typeToLayout[slideType]]
        slide = prs.slides.add_slide(layout)
        shapes = slide.shapes
        if "backgroundImage" in slideInfo:
            backgroundImage = slideInfo["backgroundImage"]
            picture = shapes.add_picture(backgroundImage, 0, 0)
            if "preserveAspectRatio" in slideInfo and slideInfo["preserveAspectRatio"]:
                width = picture.width
                height = picture.height
                aspectRatio = width / height
                if aspectRatio > 1:
                    width = prs.slide_width
                    height = int(width / aspectRatio)
                else:
                    height = prs.slide_height
                    width = int(height * aspectRatio)
                x = prs.slide_width / 2 - width / 2
                y = prs.slide_height / 2 - height / 2
                newPicture = shapes.add_picture(backgroundImage, x, y, width, height)
                shapes._spTree.remove(picture._element)
                picture = newPicture
            else:
                picture.width = prs.slide_width
                picture.height = prs.slide_height
            # Place the image in background
            shapes._spTree.remove(picture._element)
            shapes._spTree.insert(2, picture._element)

        title = shapes.title
        titleData = TextData(slideInfo["title"])
        title.text = titleData.text
        titleFont = title.text_frame.paragraphs[0].font
        titleData.applyTo(titleFont)
        if slideType == "title":
            if "subtitle" in slideInfo:
                subtitle = slide.placeholders[1]
                subtitleData = TextData(slideInfo["subtitle"])
                subtitle.text = subtitleData.text
                subtitleFont = subtitle.text_frame.paragraphs[0].font
                subtitleData.applyTo(subtitleFont)
        elif slideType == "bullets" or slideType == "content":
            if slideType == "bullets":
                body = shapes.placeholders[1]
            elif slideType == "content":
                body = shapes.placeholders[14]

            textWidth = prs.slide_width

            if "pictures" in slideInfo:
                pictures = slideInfo["pictures"]
                for picture in pictures:
                    horizontalAlignment = picture.get("horizontalAlignment", "center")
                    imageWidth, imageHeight = PIL.Image.open(picture["file"]).size
                    aspectRatio = imageHeight / imageWidth
                    width = int(prs.slide_height / aspectRatio)
                    height = prs.slide_height
                    if horizontalAlignment == "center":
                        x = prs.slide_width / 2 - width / 2
                    elif horizontalAlignment == "leading":
                        x = 0
                        textWidth -= width
                    else:  # Trailing
                        x = prs.slide_width - width
                        textWidth -= width
                    p = shapes.add_picture(picture["file"], x, 0, width, height)
                    shapes._spTree.remove(p._element)
                    shapes._spTree.insert(2, p._element)

            # body.width = textWidth
            tf = body.text_frame
            tf.clear()
            bulletData = BulletData(slideInfo)
            for bullet in bulletData.bullets:
                p = tf.add_paragraph()
                p.text = bullet.text
                bullet.applyTo(p.font)

    prs.save(description["name"])


__all__ = ["makePresentation"]
